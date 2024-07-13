import re
from pptx.enum.shapes import MSO_SHAPE_TYPE

def slide_replace_all(presentation, replacements, case_sensitive=True, whole_word=False):
    """Replace all occurrences of specified strings in the presentation.

    Args:
        presentation (Presentation): The PowerPoint presentation object.
        replacements (list): List of dictionaries with 'match' and 'replace' keys.
        case_sensitive (bool): Whether the replacements should be case-sensitive.
        whole_word (bool): Whether to match whole words only.

    Returns:
        int: Total number of replacements made.
    """
    total_replacements = 0

    # Compile regex patterns for each replacement
    patterns = []
    for rep in replacements:
        if rep.get('is_regex', False):
            pattern = rep['match']
        else:
            pattern = re.escape(rep['match'])
        
        if whole_word and not rep.get('is_regex', False):
            pattern = r'\b' + pattern + r'\b'
        
        flags = 0 if case_sensitive else re.IGNORECASE
        patterns.append((re.compile(pattern, flags), rep['replace']))

    def replace_text(text):
        nonlocal total_replacements
        for pattern, replacement in patterns:
            new_text, count = pattern.subn(replacement, text)
            total_replacements += count
            text = new_text
        return text

    def process_shape(shape):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for subshape in shape.shapes:
                process_shape(subshape)
        elif shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.text = replace_text(run.text)
        elif shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = replace_text(run.text)
        elif shape.has_chart:
            if shape.chart.has_title:
                shape.chart.chart_title.text_frame.text = replace_text(shape.chart.chart_title.text_frame.text)
            for series in shape.chart.series:
                if series.name:
                    series.name = replace_text(series.name)
            for category in shape.chart.categories:
                if category.label:
                    category.label = replace_text(category.label)

    for slide in presentation.slides:
        for shape in slide.shapes:
            process_shape(shape)

    return total_replacements
