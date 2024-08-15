import collections 
import collections.abc
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Pt, Emu

# Alignment mapping dictionary
alignment_map = {
    PP_ALIGN.LEFT: "LEFT",
    PP_ALIGN.CENTER: "CENTER",
    PP_ALIGN.RIGHT: "RIGHT",
    PP_ALIGN.JUSTIFY: "JUSTIFY",
    PP_ALIGN.DISTRIBUTE: "DISTRIBUTE",
    PP_ALIGN.JUSTIFY_LOW: "JUSTIFY_LOW",
    PP_ALIGN.THAI_DISTRIBUTE: "THAI_DISTRIBUTE"
}

def slide_add_table(presentation, slide_number, table_data, position):
    """Add a new table to a specific slide in the presentation."""
    try:
        # Get the specified slide
        slide = presentation.slides[slide_number - 1]
        
        # Extract table dimensions from data
        rows = len(table_data['data'])
        cols = len(table_data['data'][0]) if rows > 0 else 0
        
        # Create a new table
        left = Emu(position['left'])
        top = Emu(position['top'])
        width = Emu(position['width'])
        height = Emu(position['height'])
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Populate table with data and apply styles
        for row_index, row_data in enumerate(table_data['data']):
            for col_index, cell_data in enumerate(row_data):
                cell = table.cell(row_index, col_index)
                style_id, text = cell_data
                cell.text = text
                apply_style(cell, table_data['styles'][style_id - 1])
        
        # Set column widths if provided
        if 'column_widths' in table_data:
            for i, width in enumerate(table_data['column_widths']):
                table.columns[i].width = Emu(width)
        
        # Apply merged cells
        for merge_range in table_data['merged_cells']:
            start_row, start_col, end_row, end_col = merge_range
            table.cell(start_row, start_col).merge(table.cell(end_row, end_col))
        
        # Set table name
        table.name = table_data['name']
        
        return f"Table '{table_data['name']}' successfully added to slide {slide_number}"
    
    except Exception as e:
        raise ValueError(f"Error adding table to slide {slide_number}: {str(e)}")

def read_slide_table(presentation, slide_number, table_name):
    """Read a table from a specific slide and return its content in a compact JSON format."""
    slide = presentation.slides[slide_number - 1]
    table = None
    for shape in slide.shapes:
        if shape.has_table and shape.name == table_name:
            table = shape.table
            break
    
    if not table:
        raise ValueError(f"Table '{table_name}' not found on slide {slide_number}")
    
    table_data = {
        "name": table_name,
        "data": [],
        "styles": [],
        "merged_cells": [],
        "column_widths": []  # New key for column widths
    }
    
    style_map = {}
    style_counter = 1
    
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            style_id = get_or_create_style(cell, style_map, style_counter)
            row_data.append([style_id, cell.text])
        table_data["data"].append(row_data)
    
    table_data["styles"] = list(style_map.values())
    table_data["merged_cells"] = get_merged_cells(table)
    
    # Add column widths
    for column in table.columns:
        table_data["column_widths"].append(column.width.emu)
    
    return table_data

def get_or_create_style(cell, style_map, style_counter):
    """Get an existing style ID or create a new one for the given cell."""
    style = extract_cell_style(cell)
    style_key = tuple(style.items())
    
    if style_key not in style_map:
        style_map[style_key] = {"id": style_counter, **style}
        style_counter += 1
    
    return style_map[style_key]["id"]

def extract_cell_style(cell):
    """Extract the style information from a cell."""
    paragraph = cell.text_frame.paragraphs[0]
    run = paragraph.runs[0] if paragraph.runs else None
    
    style = {}
    
    # Background color
    if cell.fill.fore_color.type != None:
        bg_color = cell.fill.fore_color.rgb
        if bg_color:
            style["bg"] = rgb_to_hex(bg_color)
    
    # Font properties
    if run:
        font_props = []
        if run.font.name:
            font_props.append(run.font.name)
        if run.font.size:
            font_props.append(str(run.font.size.pt))
        font_props.append(str(run.font.bold))
        try:
            if run.font.color.rgb:
                font_props.append(rgb_to_hex(run.font.color.rgb))
        except Exception as e:
            print("Error in reading font color (may be okay if no color defined)", e)
            pass

        if font_props:
            style["font"] = ",".join(font_props)
    
    # Alignment
    if paragraph.alignment:
        style["align"] = alignment_map.get(paragraph.alignment, "LEFT")  # Default to LEFT if not found
    
    # Border
    border = extract_border_style(cell)
    if border:
        style["border"] = border
    
    return style

def extract_border_style(cell):
    """Extract the border style from a cell."""
    if not hasattr(cell._tc.tcPr, 'tcBorders'):
        return None
    
    borders = cell._tc.tcPr.tcBorders
    if not borders:
        return None
    
    border_sides = ["top", "right", "bottom", "left"]
    border_styles = []
    
    for side in border_sides:
        if hasattr(borders, side):
            border = getattr(borders, side)
            if border and border.val:
                width = border.sz.pt if border.sz else 1
                style = border.val
                color = rgb_to_hex(border.color.rgb) if border.color and border.color.rgb else "#000000"
                border_styles.append(f"{side[0]}:{width}px {style} {color}")
    
    return "; ".join(border_styles) if border_styles else None

def apply_style(cell, style):
    """Apply the given style to a cell."""
    if "bg" in style:
        cell.fill.solid()
        cell.fill.fore_color.rgb = hex_to_rgb(style["bg"])
    
    if "font" in style:
        font_props = style["font"].split(",")
        paragraph = cell.text_frame.paragraphs[0]
        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
        if len(font_props) > 0:
            run.font.name = font_props[0]
        if len(font_props) > 1:
            run.font.size = Pt(float(font_props[1]))
        if len(font_props) > 2:
            run.font.bold = font_props[2].lower() == "true"
        if len(font_props) > 3:
            run.font.color.rgb = hex_to_rgb(font_props[3])
    
    if "align" in style:
        cell.text_frame.paragraphs[0].alignment = getattr(PP_ALIGN, style["align"])
    
    if "border" in style:
        apply_border_style(cell, style["border"])

def apply_border_style(cell, border_string):    
    """Apply the border style to a cell."""
    # not working, causes error
    return
    borders = cell._tc.tcPr.get_or_add_tcBorders()
    for border_part in border_string.split(";"):
        side, style = border_part.strip().split(":")
        width, style, color = style.split()
        border = getattr(borders, side_map[side[0]])
        border.val = style
        border.sz = Pt(float(width[:-2]))  # Remove 'px' and convert to points
        border.color.rgb = hex_to_rgb(color)

def get_merged_cells(table):
    """Get the merged cell ranges in the table."""
    merged_cells = []
    for row_idx, row in enumerate(table.rows):
        for col_idx, cell in enumerate(row.cells):
            if cell.is_merge_origin:
                end_row = row_idx + cell.span_height - 1
                end_col = col_idx + cell.span_width - 1
                merged_cells.append([row_idx, col_idx, end_row, end_col])
    return merged_cells

def rgb_to_hex(rgb):
    """Convert RGB tuple to hex color string."""
    if rgb is None:
        return None
    return '#{:02x}{:02x}{:02x}'.format(*rgb)

def hex_to_rgb(hex_color):
    """Convert hex color string to RGB tuple."""
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:], 16))

side_map = {"t": "top", "r": "right", "b": "bottom", "l": "left"}


if __name__ == "__main__":
    presentation = Presentation("test.pptx")
    
    # Example usage of slide_add_table
    table_data = {
        "name": "New Table",
        "data": [
            [[1, "Header 1"], [1, "Header 2"]],
            [[2, "Data 1"], [2, "Data 2"]]
        ],
        "styles": [
            {"id": 1, "bg": "#CCCCCC", "font": "Arial,12,true,#000000", "align": "CENTER"},
            {"id": 2, "font": "Calibri,11,false,#000000", "align": "LEFT"}
        ],
        "merged_cells": [],
        "column_widths": [Emu(914400).emu, Emu(914400).emu]  # 1 inch each
    }
    
    position = {
        "left": Emu(914400).emu,  # 1 inch from left
        "top": Emu(914400).emu,   # 1 inch from top
        "width": Emu(5486400).emu, # 6 inches wide
        "height": Emu(914400).emu # 1 inch tall
    }
    
    result = slide_add_table(presentation, 1, table_data, position)
    print(result)
    
    # Example usage of read_slide_table
    table_data = read_slide_table(presentation, 1, "New Table")
    print(table_data)
    
    presentation.save("test_with_new_table.pptx")
