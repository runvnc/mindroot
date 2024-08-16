from pptx import Presentation
from pptx.chart.data import CategoryChartData
import re

def update_slide_content(context, filename, slide_number, content):
    try:
        prs = Presentation(filename)
        slide = prs.slides[slide_number - 1]
        content_applied = {key: False for key in content.keys()}  # Track which content has been applied
        print('-----------------------------------------------------')
        print(f"Updating content of slide {slide_number} in {filename}")
        num_updates = 0
        for shape in slide.shapes:
            if shape.name in content and not content_applied[shape.name]:
                print("Found shape", shape.name)
                if shape.has_text_frame:
                    update_text_frame(shape.text_frame, content[shape.name])
                    num_updates += 1
                elif shape.has_table:
                    update_table(shape.table, content[shape.name])
                    num_updates += 1
                elif shape.has_chart:
                    update_chart(shape.chart, content[shape.name])
                    num_updates += 1
                content_applied[shape.name] = True  # Mark this content as applied
        print("Done updating content")
        prs.save(filename)
        if num_updates == 0:
            return f"Error: No names matched in attempt to update slide {slide_number} in {filename}"
        return f"Updated content of slide {slide_number} in {filename}, matched {num_updates} names"
    except Exception as e:
        import traceback
        tb = traceback.format_exc()
        return f"Error updating slide content: {str(e)}\n\n{tb}"

def update_text_frame(text_frame, new_content):
    print("Updating text frame: ", text_frame.text)
    if isinstance(new_content, str):
        new_content = [new_content]
    
    for i, content in enumerate(new_content):
        if i < len(text_frame.paragraphs):
            update_paragraph(text_frame.paragraphs[i], content)
        else:
            p = text_frame.add_paragraph()
            p.text = content

def update_paragraph(paragraph, new_text):
    if len(paragraph.runs) == 1:
        # Single run: retain formatting of the first (only) run
        run = paragraph.runs[0]
        run.text = new_text
    else:
        # Multiple runs: join text, apply replacement, update first run
        whole_text = " ".join([r.text for r in paragraph.runs])
        whole_text = new_text  # Replace entire text
        
        p = paragraph._p  # the lxml element containing the `<a:p>` paragraph element
        # Remove all but the first run
        for run in paragraph.runs[1:]:
            p.remove(run._r)
        
        # Update the text of the first run
        paragraph.runs[0].text = whole_text

def update_table(table, new_content):
    if not isinstance(new_content, list) or not all(isinstance(row, list) for row in new_content):
        raise ValueError("Content for table must be a 2D list")
    
    for i, row in enumerate(new_content):
        for j, cell_content in enumerate(row):
            if i < len(table.rows) and j < len(table.columns):
                cell = table.cell(i, j)
                update_text_frame(cell.text_frame, str(cell_content))

def update_chart(chart, new_content):
    if not isinstance(new_content, dict) or 'categories' not in new_content or 'values' not in new_content:
        raise ValueError("Content for chart must be a dict with 'categories' and 'values' keys")
    
    chart_data = CategoryChartData()
    chart_data.categories = new_content['categories']
    chart_data.add_series('Series 1', new_content['values'])
    chart.replace_data(chart_data)
