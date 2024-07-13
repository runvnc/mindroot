from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.chart.data import ChartData
import json

def read_slide_content(presentation, slide_number):
    """Read and return the content of a slide with improved robustness.

    Args:
        presentation (Presentation): The PowerPoint presentation object.
        slide_number (int): The 1-based index of the slide to read.

    Returns:
        dict: A dictionary containing the slide's content and metadata.
    """
    try:
        slide = presentation.slides[slide_number - 1]
    except IndexError:
        return {"error": f"Slide number {slide_number} does not exist in the presentation."}

    content = {
        "slide_number": slide_number,
        "slide_layout": slide.slide_layout.name,
        "shapes": []
    }

    def process_shape(shape, is_in_group=False):
        shape_data = {
            "name": shape.name,
            "type": str(shape.shape_type),
            "left": shape.left,
            "top": shape.top,
            "width": shape.width,
            "height": shape.height
        }

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            shape_data["grouped_shapes"] = [process_shape(subshape, is_in_group=True) for subshape in shape.shapes]
        elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            shape_data["placeholder_type"] = str(shape.placeholder_format.type)

        if shape.has_text_frame:
            shape_data["text"] = process_text_frame(shape.text_frame)
        elif shape.has_table:
            shape_data["table"] = [[cell.text for cell in row.cells] for row in shape.table.rows]
        elif shape.has_chart:
            shape_data["chart"] = process_chart(shape.chart)
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shape_data["picture"] = {
                "filename": shape.image.filename
            }

        return shape_data

    def process_text_frame(text_frame):
        return [{
            "text": paragraph.text,
            "level": paragraph.level
        } for paragraph in text_frame.paragraphs]

    def process_chart(chart):
        chart_data = {
            "chart_type": str(chart.chart_type),
            "has_legend": chart.has_legend,
            "series": []
        }

        if chart.has_title:
            chart_data["title"] = chart.chart_title.text_frame.text

        for series in chart.series:
            series_data = {
                "name": series.name,
                "values": list(series.values)
            }
            if hasattr(series, 'categories'):
                series_data["categories"] = list(series.categories)
            chart_data["series"].append(series_data)

        return chart_data

    for shape in slide.shapes:
        content["shapes"].append(process_shape(shape))

    return content

def read_presentation_content(filename):
    """Read and return the content of all slides in a presentation.

    Args:
        filename (str): The filename of the PowerPoint presentation.

    Returns:
        dict: A dictionary containing the content of all slides.
    """
    try:
        prs = Presentation(filename)
    except Exception as e:
        return {"error": f"Failed to open presentation: {str(e)}"}

    presentation_content = {
        "slide_count": len(prs.slides),
        "slides": []
    }

    for i, slide in enumerate(prs.slides, start=1):
        slide_content = read_slide_content(prs, i)
        presentation_content["slides"].append(slide_content)

    return presentation_content

# Example usage
if __name__ == "__main__":
    filename = "example.pptx"  # Replace with your PowerPoint file
    content = read_presentation_content(filename)
    print(json.dumps(content, indent=2))
