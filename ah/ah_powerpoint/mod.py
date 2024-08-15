from ..commands import command
from pptx import Presentation, parts
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
import json
from .read_slide import read_slide_content as new_read_slide_content
from .replace_all import slide_replace_all as new_slide_replace_all
from .slide_content_updater import slide_add_table as slide_add_table_impl
from .pptx_stateless_editor import presentation_manager, extract_slide_xml, update_slide_xml, clear_slide, append_to_slide, done_appending
from .table_operations import read_slide_table as read_slide_table_impl, update_slide_table as update_slide_table_impl
import traceback

@command()
async def slide_replace_all(context, filename, slide, replacements=None, case_sensitive=True, whole_word=False):
    """Replace all occurrences of specified strings on the presentation slide.
    NOTE: for this and ALL PowerPoint commands, filenames must be specified with full absolute paths!

    Parameters:
        filename: string
        slide (int): slide number to update
        replacements: list of dicts with 'match' and 'replace' keys
        case_sensitive: boolean
        whole_word: boolean

    Example:
    { "slide_replace_all": { "filename": "/path/to/example.pptx", "slide": 1, "replacements": [{"match": "old text", "replace": "new text"}], "case_sensitive": true, "whole_word": false } }
    """
    if replacements is None:
        return "No replacements provided"

    try:
        prs = Presentation(filename)
        total_replacements = new_slide_replace_all(prs, replacements, case_sensitive, whole_word)
        prs.save(filename)
        return f"Completed {total_replacements} replacements across the presentation and saved to {filename}."
    except Exception as e:
        return f"Error during replacement: {str(e)}"

@command()
async def replace_image(context, filename, slide, image_name=None, replace_with_image_fname=None):
    """Replace an image in the presentation based on name.
    
    Parameters:
        filename: string
        slide (int): 1-based slide number
        image_name: string
        replace_with_image_fname: string (full path to new image)

    Example:
    { "replace_image": { "filename": "/path/to/presentation.pptx", "slide": 1, "image_name": "Picture 1", "replace_with_image_fname": "/absolute/path/to/new_logo.png" } }
    """
    if image_name is None or replace_with_image_fname is None:
        return "Both image_name and replace_with_image_fname must be provided"
    if slide is None:
        return "Must provide slide #"

    try:
        prs = Presentation(filename)
        replacements_count = 0

        slide = prs.slides[slide - 1]
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                if shape.name == image_name:
                    im = parts.image.Image.from_file(replace_with_image_fname)

                    slide_part, rId = shape.part, shape._element.blip_rId
                    image_part = slide_part.related_part(rId)
                    image_part.blob = im._blob

                    replacements_count += 1

        prs.save(filename)
        return f"Replaced {replacements_count} instance(s) of {image_name} with {replace_with_image_fname} in {filename}"
    except Exception as e:
        return f"Error replacing image: {str(e)}"

@command()
async def read_slide_content(context, filename, slide_number):
    """Read and return the content of a slide.
    Example:
    { "read_slide_content": { "filename": "/path/to/presentation.pptx", "slide_number": 1 } }
    """
    try:
        prs = Presentation(filename)
        content = new_read_slide_content(prs, slide_number)
        return json.dumps(content)
    except Exception as e:
        return json.dumps({"error": str(e)})

@command()
async def update_slide_content(context, filename, slide_number, content):
    """Update a slide with content.
    For lists of text in groups or text that can be updated with search and replace, use
    slide_replace_all instead. Text frame updates can be string or list of strings.
    
    For tables, the number of rows in the provided data must match the existing table.
    Example:
    { "update_slide_content": { "filename": "/path/to/presentation.pptx", "slide_number": 1, "content": {"title": "New Title", "subtitle": "New Subtitle", "Table 2": [["Header 1", "Header 2"], ["Row 1 Col 1", "Row 1 Col 2"]]} } }
    """
    try:
        return new_update_slide_content(context, filename, slide_number, content)
    except Exception as e:
        return f"Error updating slide content: {str(e)}"

@command()
async def create_chart(context, filename, slide_number, chart_type, data, position):
    """Create a chart on the specified slide.
    Example:
    { "create_chart": { "filename": "/path/to/presentation.pptx", "slide_number": 1, "chart_type": "BAR_CLUSTERED", "data": {"categories": ["A", "B", "C"], "values": [1, 2, 3]}, "position": {"left": 1, "top": 2, "width": 8, "height": 5} } }
    """
    try:
        prs = Presentation(filename)
        slide = prs.slides[slide_number - 1]
        chart_data = CategoryChartData()
        chart_data.categories = data['categories']
        chart_data.add_series('Series 1', data['values'])
        x, y, cx, cy = [Inches(position[key]) for key in ['left', 'top', 'width', 'height']]
        chart = slide.shapes.add_chart(
            getattr(XL_CHART_TYPE, chart_type), x, y, cx, cy, chart_data
        ).chart
        prs.save(filename)
        return f"Created {chart_type} chart on slide {slide_number} in {filename}"
    except Exception as e:
        return f"Error creating chart: {str(e)}"

@command()
async def open_presentation_for_xml_edit(context, filename):
    """Open a presentation for XML editing.
    NOTE: This only appplies to xml presentation commands!\

    Example:
    { "open_presentation_for_xml_edit": { "filename": "/path/to/example.pptx" } }
    """
    try:
        result = presentation_manager.open_presentation(filename)
        return f"Presentation {filename} opened for XML editing" if result else "Failed to open presentation"
    except Exception as e:
        return f"Error opening presentation: {str(e)}"

@command()
async def close_presentation_after_xml_edit(context, filename):
    """Close a presentation after XML editing.
    NOTE: This only appplies to xml presentation commands!

    Example:
    { "close_presentation_after_xml_edit": { "filename": "/path/to/example.pptx" } }
    """
    try:
        result = presentation_manager.close_presentation(filename)
        return f"Presentation {filename} closed after XML editing" if result else "Failed to close presentation"
    except Exception as e:
        return f"Error closing presentation: {str(e)}"

@command()
async def extract_slide_xml_content(context, filename, slide_number):
    """Extract XML content from a slide.
    Example:
    { "extract_slide_xml_content": { "filename": "/path/to/example.pptx", "slide_number": 1 } }
    """
    try:
        xml_content = extract_slide_xml(filename, slide_number)
        return xml_content
    except Exception as e:
        return f"Error extracting slide XML: {str(e)}"

@command()
async def update_slide_xml_content(context, filename, slide_number, new_xml):
    """Update XML content of a slide.
    Example:
    { "update_slide_xml_content": { "filename": "/path/to/example.pptx", "slide_number": 1, "new_xml": "<p:sld>...</p:sld>" } }
    """
    try:
        result = update_slide_xml(filename, slide_number, new_xml)
        return f"Updated XML content for slide {slide_number} in {filename}" if result else "Failed to update slide XML"
    except Exception as e:
        return f"Error updating slide XML: {str(e)}"

@command()
async def append_to_slide_xml_content(context, filename, slide_number, xml_fragment):
    """Append XML fragment to a slide's content.
       Use this if you aren't sure you can output the full slide contents in one go.
       If you just cleared, don't forget "header"/namespace tags etc.
    Example:
    { "append_to_slide_xml_content": { "filename": "/path/to/example.pptx", "slide_number": 1, "xml_fragment": "<p:sp>...</p:sp>" } }
    """
    try:
        result = append_to_slide(filename, slide_number, xml_fragment)
        # print out in magenta for debugging
        print(f"appended: \033[35m{xml_fragment}\033[0m")
        return f"Appended XML fragment to slide {slide_number} in {filename}" if result else "Failed to append to slide XML"
    except Exception as e:
        return f"Error appending to slide XML: {str(e)}"

@command()
async def clear_slide_xml_content(context, filename, slide_number):
    """Clear all XML content from a slide.
    NOTE: Only use this with other xml commands.

    Note that this includes ALL XML, including namespaces and other attributes.

    Example:
    { "clear_slide_xml_content": { "filename": "/path/to/example.pptx", "slide_number": 1 } }
    """
    try:
        result = clear_slide(filename, slide_number)
        return f"Cleared XML content from slide {slide_number} in {filename}" if result else "Failed to clear slide XML"
    except Exception as e:
        return f"Error clearing slide XML: {str(e)}"

@command()
async def done_appending_slide(context, filename, slide_number):
    """Use this when you are finished appending XML to a slide.
    NOTE: Only use this with other xml commands.

    Example:
    { "done_appending_slide": { "filename": "/path/to/example.pptx", "slide_number": 1 } }
    """
    try:
        result = done_appending(filename, slide_number)
        return f"Finished appending. Presentation saved."
    except Exception as e:
        return f"Error finishing appending to slide XML: {str(e)}"
 
@command()
async def save_presentation_after_xml_edit(context, filename):
    """Save the presentation after XML editing.
    NOTE: This only appplies to xml presentation commands!

    Example:
    { "save_presentation_after_xml_edit": { "filename": "/path/to/example.pptx" } }
    """
    try:
        result = presentation_manager.save_presentation(filename)
        return f"Saved presentation {filename} after XML editing" if result else "Failed to save presentation"
    except Exception as e:
        return f"Error saving presentation: {str(e)}"

@command()
async def read_slide_table(context, filename, slide_number, table_name):
    """Read a table from a specific slide and return its content in a compact JSON format.

    Parameters:
        filename (str): Full path to the PowerPoint file.
        slide_number (int): 1-based slide number containing the table.
        table_name (str): Name of the table shape to read.

    Returns:
        str: A JSON string representing the table structure, content, styling, and column widths.

    The returned JSON structure includes:
    {
        "name": "Table Name",
        "data": [[[cell_data], [cell_data]],[ ] ...],
        "styles": [style_data],
        "merged_cells": [[start_row, start_col, end_row, end_col], ...],
        "column_widths": [width1, width2, ...]
    }

    Where:
    - data is list of rows, each row is a list of cell_data
    - cell_data is [style_id, "cell content"]
    - style_data is {"id": style_id, "bg": "#RRGGBB", "font": "name,size,bold,color", "align": "alignment", "border": "border_style"}
    - column_widths are in EMU (English Metric Units)

    Example:
    { "read_slide_table": { "filename": "/path/to/presentation.pptx", "slide_number": 1, "table_name": "Table 1" } }
    """
    try:
        prs = Presentation(filename)
        table_data = read_slide_table_impl(prs, slide_number, table_name)
        return json.dumps(table_data)
    except Exception as e:
        return json.dumps({"error": str(e)})

@command()
async def slide_add_table(context, filename, slide_number, table_data, position):
    """Add a new table to a specific slide in the presentation.
    Note: the presentation will automatically be saved after adding the table.

    Parameters:
        filename (str): Full path to the PowerPoint file.
        slide_number (int): 1-based slide number where the table will be added.
        table_data (dict): Dict representing the table structure, content, and styling.
        position (dict): Dict specifying the position and size of the table on the slide.

    Returns:
        str: A message indicating success or describing an error.

    The table_data should have the following structure:
    {
        "name": "Table Name",
        "data": [[[cell_data], [cell_data], ...], [[cell_data], [cell_data], ...], ...],
        "styles": [style_data],
        "merged_cells": [[start_row, start_col, end_row, end_col], ...],
        "column_widths": [width1, width2, ...] (optional)
    }

    Where:
    - cell_data is [style_id, "cell content"]
    - style_data is {"id": style_id, "bg": "#RRGGBB", "font": "name,size,bold,color", "align": "ALIGNMENT"}
    - column_widths are in EMU (English Metric Units)

    The position dict should have the following structure:
    {
        "left": left_position,
        "top": top_position,
        "width": table_width,
        "height": table_height
    }
    All values in the position dict should be in EMU (English Metric Units).

    Example:
    {
        "slide_add_table": {
            "filename": "/path/to/presentation.pptx",
            "slide_number": 1,
            "table_data": {
                "name": "Sample Table",
                "data": [
                    [[1, "Header 1"], [1, "Header 2"]],
                    [[2, "Data 1"], [2, "Data 2"]]
                ],
                "styles": [
                    {"id": 1, "bg": "#CCCCCC", "font": "Arial,12,true,#000000", "align": "CENTER"},
                    {"id": 2, "font": "Calibri,11,false,#000000", "align": "LEFT"}
                ],
                "merged_cells": [],
                "column_widths": [914400, 914400]  // Optional: 1 inch each (914400 EMU = 1 inch)
            },
            "position": {
                "left": 1828800,  // 2 inches from left
                "top": 1828800,   // 2 inches from top
                "width": 5486400, // 6 inches wide
                "height": 1828800 // 2 inches tall
            }
        }
    }
    """
    try:
        prs = Presentation(filename)
        result = slide_add_table_impl(prs, slide_number, table_data, position)
        prs.save(filename)
        return result
    except Exception as e:
        return json.dumps({"error": str(e)})


