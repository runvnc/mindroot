from ..commands import command
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
import json
import io

# Module-level cache for the presentation
_presentation_cache = None

def _get_presentation():
    global _presentation_cache
    if _presentation_cache is None:
        raise ValueError("No presentation is currently cached. Use open_presentation first.")
    return Presentation(io.BytesIO(_presentation_cache))

def _save_presentation_cache(prs):
    global _presentation_cache
    buffer = io.BytesIO()
    prs.save(buffer)
    _presentation_cache = buffer.getvalue()

@command()
async def open_presentation(filename, context=None):
    """Open a PowerPoint presentation and cache it in memory.
    Example:
    { "open_presentation": { "filename": "example.pptx" } }
    """
    prs = Presentation(filename)
    _save_presentation_cache(prs)
    context.data['current_presentation'] = filename
    slides_info = [{'number': i+1, 'layout': slide.slide_layout.name} for i, slide in enumerate(prs.slides)]
    return f"Opened and cached presentation {filename}. Slides: {json.dumps(slides_info)}"

@command()
async def list_placeholders(slide_number, context=None):
    """List all placeholders in a specific slide.
    Example:
    { "list_placeholders": { "slide_number": 1 } }
    """
    prs = _get_presentation()
    slide = prs.slides[slide_number - 1]
    placeholders = {shape.placeholder_format.idx: shape.name for shape in slide.shapes if shape.is_placeholder}
    return json.dumps(placeholders)

@command()
async def fill_placeholder(slide_number, placeholder_index, content, context=None):
    """Fill a specific placeholder with content.
    Example:
    { "fill_placeholder": { "slide_number": 1, "placeholder_index": 0, "content": "Hello, World!" } }
    """
    prs = _get_presentation()
    slide = prs.slides[slide_number - 1]
    shape = slide.placeholders[placeholder_index]
    
    if shape.has_text_frame:
        shape.text_frame.text = content
    elif shape.has_table:
        if isinstance(content, list) and all(isinstance(row, list) for row in content):
            for i, row in enumerate(content):
                for j, cell_content in enumerate(row):
                    if i < len(shape.table.rows) and j < len(shape.table.columns):
                        shape.table.cell(i, j).text = str(cell_content)
        else:
            raise ValueError("Content for table must be a 2D list")
    elif shape.has_chart:
        if isinstance(content, dict) and 'categories' in content and 'values' in content:
            chart_data = CategoryChartData()
            chart_data.categories = content['categories']
            chart_data.add_series('Series 1', content['values'])
            shape.chart.replace_data(chart_data)
        else:
            raise ValueError("Content for chart must be a dict with 'categories' and 'values' keys")
    
    _save_presentation_cache(prs)
    return f"Filled placeholder {placeholder_index} on slide {slide_number}"

@command()
async def add_slide(layout_name, context=None):
    """Add a new slide with the specified layout.
    Example:
    { "add_slide": { "layout_name": "Title Slide" } }
    """
    prs = _get_presentation()
    layout = next((layout for layout in prs.slide_layouts if layout.name == layout_name), None)
    if layout is None:
        return f"Layout '{layout_name}' not found"
    new_slide = prs.slides.add_slide(layout)
    _save_presentation_cache(prs)
    return f"Added new slide with layout '{layout_name}'. Slide number: {len(prs.slides)}"

@command()
async def delete_slide(slide_number, context=None):
    """Delete the specified slide.
    Example:
    { "delete_slide": { "slide_number": 2 } }
    """
    prs = _get_presentation()
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[slide_number - 1])
    _save_presentation_cache(prs)
    return f"Deleted slide number {slide_number}"

@command()
async def read_slide_content(slide_number, context=None):
    """Read and return the content of a slide.
    Example:
    { "read_slide_content": { "slide_number": 1 } }
    """
    prs = _get_presentation()
    slide = prs.slides[slide_number - 1]
    content = {}
    for shape in slide.shapes:
        if shape.has_text_frame:
            content[shape.name] = shape.text_frame.text
        elif shape.has_table:
            content[shape.name] = [[cell.text for cell in row.cells] for row in shape.table.rows]
        elif shape.has_chart:
            content[shape.name] = "Chart: " + shape.chart.chart_type
    return json.dumps(content)

@command()
async def update_slide_content(slide_number, content_json, context=None):
    """Update a slide with content provided in JSON format.
    Example:
    { "update_slide_content": { "slide_number": 1, "content_json": {"title": "New Title", "subtitle": "New Subtitle"} } }
    """
    prs = _get_presentation()
    slide = prs.slides[slide_number - 1]
    content = json.loads(content_json)
    for shape in slide.shapes:
        if shape.name in content:
            if shape.has_text_frame:
                shape.text_frame.text = content[shape.name]
            elif shape.has_table:
                if isinstance(content[shape.name], list) and all(isinstance(row, list) for row in content[shape.name]):
                    for i, row in enumerate(content[shape.name]):
                        for j, cell_content in enumerate(row):
                            if i < len(shape.table.rows) and j < len(shape.table.columns):
                                shape.table.cell(i, j).text = str(cell_content)
                else:
                    raise ValueError(f"Content for table '{shape.name}' must be a 2D list")
            elif shape.has_chart:
                if isinstance(content[shape.name], dict) and 'categories' in content[shape.name] and 'values' in content[shape.name]:
                    chart_data = CategoryChartData()
                    chart_data.categories = content[shape.name]['categories']
                    chart_data.add_series('Series 1', content[shape.name]['values'])
                    shape.chart.replace_data(chart_data)
                else:
                    raise ValueError(f"Content for chart '{shape.name}' must be a dict with 'categories' and 'values' keys")
    _save_presentation_cache(prs)
    return f"Updated content of slide {slide_number}"

@command()
async def create_chart(slide_number, chart_type, data, position, context=None):
    """Create a chart on the specified slide.
    Example:
    { "create_chart": { "slide_number": 1, "chart_type": "BAR_CLUSTERED", "data": {"categories": ["A", "B", "C"], "values": [1, 2, 3]}, "position": {"left": 1, "top": 2, "width": 8, "height": 5} } }
    """
    prs = _get_presentation()
    slide = prs.slides[slide_number - 1]
    chart_data = CategoryChartData()
    chart_data.categories = data['categories']
    chart_data.add_series('Series 1', data['values'])
    x, y, cx, cy = [Inches(position[key]) for key in ['left', 'top', 'width', 'height']]
    chart = slide.shapes.add_chart(
        getattr(XL_CHART_TYPE, chart_type), x, y, cx, cy, chart_data
    ).chart
    _save_presentation_cache(prs)
    return f"Created {chart_type} chart on slide {slide_number}"

@command()
async def list_layouts(context=None):
    """List all available slide layouts in the presentation.
    Example:
    { "list_layouts": {} }
    """
    prs = _get_presentation()
    layouts = [layout.name for layout in prs.slide_layouts]
    return json.dumps(layouts)

@command()
async def save_presentation(filename=None, context=None):
    """Save the presentation to disk, optionally with a new filename.
    Example:
    { "save_presentation": { "filename": "updated_presentation.pptx" } }
    """
    prs = _get_presentation()
    if filename is None:
        filename = context.data['current_presentation']
    prs.save(filename)
    context.data['current_presentation'] = filename
    return f"Saved presentation as {filename}"
